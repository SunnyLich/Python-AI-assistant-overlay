"""
core/llm.py — Cloud LLM client with streaming support.

Supports:
  - Groq (OpenAI-compatible, fast TTFT)
  - OpenAI
  - Anthropic Claude

Clients are module-level singletons so the TLS connection is reused across
calls, eliminating handshake overhead from every request.
"""
from __future__ import annotations
import config
from pathlib import Path
from core import secret_store
from core.tool_registry import ToolRegistry, ToolSpec
from typing import Generator

_TOOL_REGISTRY = ToolRegistry()


def _log_context(
    reason: str,
    text: str,
    max_line: int = 120,
    max_lines: int = 12,
    max_chars: int = 1200,
) -> None:
    """Print a compact preview of a context block for debugging."""
    import time

    ts = time.strftime("%H:%M:%S")

    def _trim(line: str) -> str:
        return line if len(line) <= max_line else line[:max_line] + "…"

    lines = [_trim(l) for l in text.splitlines() if l.strip()]
    truncated = False

    if len(lines) > max_lines:
        lines = lines[:max_lines]
        truncated = True

    body = "\n  ".join(lines) if lines else "[empty]"
    if len(body) > max_chars:
        body = body[:max_chars].rstrip() + "…"
        truncated = True
    if truncated and body != "[empty]":
        body += "\n  [preview truncated]"

    print(f"[llm {ts}] Context — {reason}:\n  {body}")


def _ambient_document_max_chars() -> int:
    return config.CONTEXT_AMBIENT_DOCUMENT_MAX_CHARS


def _tool_document_max_chars() -> int:
    return config.CONTEXT_TOOL_DOCUMENT_MAX_CHARS


def _read_document_file(path: str, max_chars: int | None = None) -> str:
    """Read a local document file and return its plain text."""
    import os
    if max_chars is None:
        max_chars = _ambient_document_max_chars()
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".docx":
            from docx import Document  # type: ignore
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext in (".xlsx", ".xls"):
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        parts.append("\t".join(cells))
            text = "\n".join(parts)
        elif ext == ".pptx":
            from pptx import Presentation  # type: ignore
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                parts.append(line)
            text = "\n".join(parts)
        elif ext == ".pdf":
            import pypdf  # type: ignore
            reader = pypdf.PdfReader(path)
            parts = []
            for i, page in enumerate(reader.pages, 1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    parts.append(f"[Page {i}]\n{page_text.strip()}")
            text = "\n\n".join(parts)
        elif ext in (".odt", ".ods", ".odp"):
            from odf import text as odf_text, teletype  # type: ignore
            from odf.opendocument import load as odf_load  # type: ignore
            doc = odf_load(path)
            paragraphs = doc.getElementsByType(odf_text.P)
            text = "\n".join(
                teletype.extractText(p) for p in paragraphs
                if teletype.extractText(p).strip()
            )
        elif ext in (".txt", ".md", ".csv", ".py", ".js", ".ts",
                     ".json", ".xml", ".html", ".log"):
            with open(path, encoding="utf-8", errors="replace") as f:
                text = f.read()
        else:
            return f"File type {ext!r} is not supported for reading."
        if len(text) > max_chars:
            text = text[:max_chars] + "\n[…truncated]"
        # Redact sensitive data before the text reaches the LLM.
        from core.context_fetcher import _redact  # noqa: PLC0415
        text = _redact(text)
        _log_context(f"tool: read_document — read {path!r}", text)
        return text
    except Exception as e:
        return f"Failed to read {path!r}: {e}"


def _read_document_paths(
    paths: list[str],
    max_chars_per_doc: int | None = None,
) -> str:
    """Read multiple local document files and join readable results."""
    import os

    if max_chars_per_doc is None:
        max_chars_per_doc = _ambient_document_max_chars()
    parts: list[str] = []
    for path in paths:
        text = _read_document_file(path, max_chars=max_chars_per_doc)
        if text and not text.startswith(("Could not", "File type", "Failed to")):
            parts.append(f"[{os.path.basename(path)}]\n{text}")
    return "\n\n".join(parts)


def _execute_get_context(inputs: dict) -> str:
    """Built-in context tool: fetch a browser page or open document text."""
    url = inputs.get("url", "").strip()
    if url:
        from core.context_fetcher import fetch_browser_content_for_tool
        result = fetch_browser_content_for_tool(url)
        _log_context(
            f"tool: get_context (browser) — {url!r}",
            result or "",
        )
        return result or f"Could not fetch content from {url!r}."

    from core.context_fetcher import get_all_open_document_paths

    paths = get_all_open_document_paths()
    if not paths:
        return "Could not determine any open document paths from supported app windows."

    text = _read_document_paths(paths, max_chars_per_doc=_tool_document_max_chars())
    if text:
        return text
    return "Open document windows were detected, but none of their file types were readable."


def _execute_git_status(inputs: dict) -> str:
    cwd = inputs.get("cwd") or config.TOOL_GIT_ROOT
    return _run_read_only_command(["git", "status", "--short"], cwd=cwd)


def _execute_git_diff(inputs: dict) -> str:
    cwd = inputs.get("cwd") or config.TOOL_GIT_ROOT
    return _run_read_only_command(["git", "diff", "--", "."], cwd=cwd)


def _run_read_only_command(args: list[str], cwd: str) -> str:
    import subprocess
    from pathlib import Path

    root = Path(cwd or ".").expanduser().resolve()
    completed = subprocess.run(
        args,
        cwd=str(root),
        capture_output=True,
        text=True,
        timeout=20,
    )
    output = (completed.stdout or completed.stderr or "").strip()
    return output[:12000] or f"{' '.join(args)} returned no output."


def _execute_github_repo(inputs: dict) -> str:
    repo = str(inputs.get("repo") or "").strip()
    if not repo:
        return "Missing repo. Use owner/name."
    return _github_get_json(f"https://api.github.com/repos/{repo}")


def _execute_github_issue(inputs: dict) -> str:
    repo = str(inputs.get("repo") or "").strip()
    number = str(inputs.get("number") or "").strip()
    if not repo or not number:
        return "Missing repo or number."
    return _github_get_json(f"https://api.github.com/repos/{repo}/issues/{number}")


def _github_get_json(url: str) -> str:
    import json
    import urllib.request
    from core import github_auth

    token = github_auth.get_valid_access_token()
    if not token:
        return "GitHub OAuth is not configured. Sign in from Settings."
    req = urllib.request.Request(
        url,
        headers={
            "Accept": "application/vnd.github+json",
            "Authorization": f"Bearer {token}",
            "User-Agent": "python-ai-overlay",
            "X-GitHub-Api-Version": "2022-11-28",
        },
    )
    with urllib.request.urlopen(req, timeout=20) as resp:
        data = json.loads(resp.read())
    return json.dumps(data, indent=2, ensure_ascii=False)[:12000]


def _register_builtin_tools() -> None:
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="web_search",
            description="Search the web for current information.",
            input_schema={"type": "object", "properties": {}, "required": []},
            server_schema={
                "type": "web_search_20250305",
                "name": "web_search",
                "max_uses": 2,
            },
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="get_context",
            description=(
                "Retrieve additional context the user can see. "
                "Pass a URL to fetch a web page; omit it to read open local "
                "documents from supported apps "
                "(Word, Excel, PowerPoint, PDF, LibreOffice, Notepad, etc.)."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "url": {
                        "type": "string",
                        "description": (
                            "A web page URL (http:// or https://) to fetch. "
                            "Omit this field to read open local documents instead."
                        ),
                    }
                },
                "required": [],
            },
            executor=_execute_get_context,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="git_status",
            description="Return read-only git status for the configured local repository.",
            input_schema={
                "type": "object",
                "properties": {"cwd": {"type": "string", "description": "Optional local repo path."}},
                "required": [],
            },
            executor=_execute_git_status,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="git_diff",
            description="Return read-only git diff for the configured local repository.",
            input_schema={
                "type": "object",
                "properties": {"cwd": {"type": "string", "description": "Optional local repo path."}},
                "required": [],
            },
            executor=_execute_git_diff,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="github_repo",
            description="Fetch GitHub repository metadata using the signed-in GitHub OAuth account.",
            input_schema={
                "type": "object",
                "properties": {"repo": {"type": "string", "description": "Repository in owner/name form."}},
                "required": ["repo"],
            },
            executor=_execute_github_repo,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="github_issue",
            description="Fetch a GitHub issue or pull request by repository and number.",
            input_schema={
                "type": "object",
                "properties": {
                    "repo": {"type": "string", "description": "Repository in owner/name form."},
                    "number": {"type": "integer", "description": "Issue or pull request number."},
                },
                "required": ["repo", "number"],
            },
            executor=_execute_github_issue,
        )
    )


def _get_tool_schemas() -> list[dict]:
    return _TOOL_REGISTRY.schemas(include_server_tools=True)


def _execute_model_tool(name: str, inputs: dict) -> str:
    return _TOOL_REGISTRY.execute(name, inputs)


_register_builtin_tools()


def read_active_document_for_context() -> str:
    """
    Read all open doc-app windows (foreground and background) and return their
    redacted plain text for proactive injection into the system prompt.
    Multiple documents are separated by per-file headers.
    Returns "" if no readable documents are found.
    """
    from core.context_fetcher import get_all_open_document_paths

    paths = get_all_open_document_paths()
    if not paths:
        return ""
    return _read_document_paths(paths)


# ------------------------------------------------------------------
# Singleton clients — initialised once, reused across all requests
# ------------------------------------------------------------------
_openai_client = None
_anthropic_client = None
_chat_openai_client = None
_chat_anthropic_client = None
_vision_openai_client = None
_vision_anthropic_client = None
_codex_client = None
_chat_codex_client = None


def reset_clients() -> None:
    """Discard all cached API clients so they are rebuilt with the current config."""
    global _openai_client, _anthropic_client, _chat_openai_client, _chat_anthropic_client
    global _vision_openai_client, _vision_anthropic_client, _codex_client, _chat_codex_client
    _openai_client = _anthropic_client = None
    _chat_openai_client = _chat_anthropic_client = None
    _vision_openai_client = _vision_anthropic_client = None
    _codex_client = _chat_codex_client = None
    _TOOL_REGISTRY.plugin_dir = Path(config.TOOL_PLUGIN_DIR)
    _TOOL_REGISTRY.refresh()


# ------------------------------------------------------------------
# Codex (ChatGPT subscription) — custom httpx transport + client
# ------------------------------------------------------------------

class _CodexTransport:
    """
    httpx-compatible BaseTransport that:
      • strips the placeholder API-key authorization header
      • injects ``Authorization: Bearer <access_token>``
      • adds ``ChatGPT-Account-Id`` and ``originator`` headers
    Tokens are fetched (and refreshed) lazily on every request.
    """

    def __init__(self):
        import httpx
        self._inner = httpx.HTTPTransport()

    def handle_request(self, request):
        import httpx
        from core import chatgpt_auth

        token      = chatgpt_auth.get_valid_access_token()
        account_id = chatgpt_auth.get_account_id()

        # Rebuild raw headers, dropping the dummy API-key auth header
        raw: list[tuple[bytes, bytes]] = [
            (name, value)
            for name, value in request.headers.raw
            if name.lower() != b"authorization"
        ]
        if token:
            raw.append((b"authorization", f"Bearer {token}".encode()))
        if account_id:
            raw.append((b"chatgpt-account-id", account_id.encode()))
        raw.append((b"originator", b"opencode"))

        new_req = httpx.Request(
            method=request.method,
            url=request.url,
            headers=raw,
            content=request.content,
            extensions=request.extensions,
        )
        return self._inner.handle_request(new_req)

    def close(self):
        self._inner.close()


def _get_codex_client():
    global _codex_client
    if _codex_client is None:
        import httpx
        from openai import OpenAI
        _codex_client = OpenAI(
            api_key="chatgpt-oauth-dummy",
            base_url="https://chatgpt.com/backend-api/codex",
            http_client=httpx.Client(transport=_CodexTransport()),
        )
    return _codex_client


def _get_chat_codex_client():
    """Returns the same singleton as _get_codex_client() — same endpoint."""
    return _get_codex_client()


# ------------------------------------------------------------------
# Config sanity checks — raise early with actionable messages
# ------------------------------------------------------------------

def _api_key_for(provider: str) -> str:
    p = provider.lower()
    if p in ("groq",):
        return config.GROQ_API_KEY
    if p == "openai":
        return config.OPENAI_API_KEY
    if p == "anthropic":
        return config.ANTHROPIC_API_KEY
    if p == "chatgpt":
        return "chatgpt-oauth"   # no static key — auth is via OAuth tokens
    if p == "copilot":
        return "copilot-token"    # stored in OS keychain
    return ""


def _credential_source_for_provider(provider: str) -> str:
    p = provider.lower()
    if p == "groq":
        return secret_store.secret_source("GROQ_API_KEY")
    if p == "openai":
        return secret_store.secret_source("OPENAI_API_KEY")
    if p == "anthropic":
        return secret_store.secret_source("ANTHROPIC_API_KEY")
    if p == "chatgpt":
        return "chatgpt-oauth"
    if p == "copilot":
        return "copilot-keychain"
    return "none"


def _log_model_route(kind: str, provider: str, model: str, use_tools: bool = False) -> None:
    import time

    ts = time.strftime("%H:%M:%S")
    tool_note = " tools=on" if use_tools else ""
    print(
        f"[llm {ts}] Route ({kind}): provider={provider or '[unset]'} "
        f"model={model or '[unset]'} credential={_credential_source_for_provider(provider)}{tool_note}"
    )


def _parse_model_fallbacks(raw: str) -> list[tuple[str, str]]:
    """Parse provider:model fallback lines from settings."""
    routes: list[tuple[str, str]] = []
    for part in raw.replace(";", "\n").splitlines():
        item = part.strip()
        if not item or item.startswith("#"):
            continue
        if ":" not in item:
            continue
        provider, model = item.split(":", 1)
        provider = provider.strip().lower()
        model = model.strip()
        if provider and model:
            routes.append((provider, model))
    return routes


def _route_candidates(provider: str, model: str, fallback_raw: str) -> list[tuple[str, str]]:
    routes: list[tuple[str, str]] = []
    if provider and model:
        routes.append((provider.lower(), model))
    for candidate in _parse_model_fallbacks(fallback_raw):
        if candidate not in routes:
            routes.append(candidate)
    return routes


def _dynamic_openai_client(provider: str):
    from openai import OpenAI

    if provider == "groq":
        return OpenAI(api_key=config.GROQ_API_KEY, base_url="https://api.groq.com/openai/v1")
    return OpenAI(api_key=config.OPENAI_API_KEY)


def _dynamic_anthropic_client():
    import anthropic

    return anthropic.Anthropic(api_key=config.ANTHROPIC_API_KEY)


_CHATGPT_SUPPORTED_MODELS = {
    "gpt-5.5", "gpt-5.4", "gpt-5.4-mini",
    "gpt-5.3-codex", "gpt-5.3-codex-spark", "gpt-5.2",
}


def _check_llm_config() -> None:
    if not config.LLM_MODEL:
        raise ValueError(
            "LLM_MODEL is not set in .env. "
            "Add LLM_MODEL=<model> (e.g. llama3-8b-8192 for Groq)."
        )
    if config.LLM_PROVIDER.lower() == "chatgpt":
        from core import chatgpt_auth
        if not chatgpt_auth.get_tokens():
            raise ValueError(
                "LLM_PROVIDER is set to 'chatgpt' but you are not logged in. "
                "Open Settings → LLM and sign in with your ChatGPT account."
            )
        if config.LLM_MODEL.lower() not in _CHATGPT_SUPPORTED_MODELS:
            raise ValueError(
                f"Model '{config.LLM_MODEL}' is not supported by the ChatGPT Codex endpoint. "
                f"Use one of: {', '.join(sorted(_CHATGPT_SUPPORTED_MODELS))}"
            )
        return
    if config.LLM_PROVIDER.lower() == "copilot":
        from core import copilot_auth
        if not copilot_auth.get_token():
            raise ValueError(
                "LLM_PROVIDER is set to 'copilot' but no GitHub Copilot token is stored. "
                "Open Settings -> LLM and save a GitHub Copilot token."
            )
        return
    if not _api_key_for(config.LLM_PROVIDER):
        raise ValueError(
            f"API key for LLM_PROVIDER='{config.LLM_PROVIDER}' is not configured. "
            "Add the matching key in Settings so it can be stored in the OS keychain."
        )


def _check_chat_llm_config() -> None:
    if not config.CHAT_LLM_MODEL:
        raise ValueError(
            "CHAT_LLM_MODEL is not set in .env. "
            "Add CHAT_LLM_MODEL=<model> or leave unset to inherit LLM_MODEL."
        )
    if config.CHAT_LLM_PROVIDER.lower() == "chatgpt":
        from core import chatgpt_auth
        if not chatgpt_auth.get_tokens():
            raise ValueError(
                "CHAT_LLM_PROVIDER is set to 'chatgpt' but you are not logged in. "
                "Open Settings → LLM and sign in with your ChatGPT account."
            )
        if config.CHAT_LLM_MODEL.lower() not in _CHATGPT_SUPPORTED_MODELS:
            raise ValueError(
                f"Model '{config.CHAT_LLM_MODEL}' is not supported by the ChatGPT Codex endpoint. "
                f"Use one of: {', '.join(sorted(_CHATGPT_SUPPORTED_MODELS))}"
            )
        return
    if config.CHAT_LLM_PROVIDER.lower() == "copilot":
        from core import copilot_auth
        if not copilot_auth.get_token():
            raise ValueError(
                "CHAT_LLM_PROVIDER is set to 'copilot' but no GitHub Copilot token is stored. "
                "Open Settings -> LLM and save a GitHub Copilot token."
            )
        return
    if not _api_key_for(config.CHAT_LLM_PROVIDER):
        raise ValueError(
            f"API key for CHAT_LLM_PROVIDER='{config.CHAT_LLM_PROVIDER}' is not configured. "
            "Add the matching key in Settings so it can be stored in the OS keychain."
        )


def _check_vision_config() -> None:
    if not config.VISION_LLM_MODEL:
        raise ValueError(
            "VISION_LLM_MODEL is not set in .env. "
            "Screen-snip queries require a vision-capable model. "
            "Example: VISION_LLM_PROVIDER=anthropic  VISION_LLM_MODEL=claude-opus-4-5"
        )
    if not config.VISION_LLM_PROVIDER:
        raise ValueError(
            "VISION_LLM_PROVIDER is not set in .env. "
            "Add VISION_LLM_PROVIDER=anthropic (or openai or chatgpt) alongside VISION_LLM_MODEL."
        )
    if config.VISION_LLM_PROVIDER.lower() == "chatgpt":
        from core import chatgpt_auth
        if not chatgpt_auth.get_tokens():
            raise ValueError(
                "VISION_LLM_PROVIDER is set to 'chatgpt' but you are not logged in. "
                "Open Settings → LLM and sign in with your ChatGPT account."
            )
        return
    if not _api_key_for(config.VISION_LLM_PROVIDER):
        raise ValueError(
            f"API key for VISION_LLM_PROVIDER='{config.VISION_LLM_PROVIDER}' is not configured. "
            "Add the matching key in Settings so it can be stored in the OS keychain."
        )


def _check_route_config(provider: str, model: str, route_name: str) -> None:
    if not provider or not model:
        raise ValueError(f"{route_name} route is missing provider or model.")
    if provider == "chatgpt":
        from core import chatgpt_auth
        if not chatgpt_auth.get_tokens():
            raise ValueError(f"{route_name} route uses chatgpt but you are not logged in.")
        if model.lower() not in _CHATGPT_SUPPORTED_MODELS:
            raise ValueError(
                f"Model '{model}' is not supported by the ChatGPT Codex endpoint. "
                f"Use one of: {', '.join(sorted(_CHATGPT_SUPPORTED_MODELS))}"
            )
        return
    if provider == "copilot":
        from core import copilot_auth
        if not copilot_auth.get_token():
            raise ValueError(f"{route_name} route uses copilot but no GitHub Copilot token is stored.")
        return
    if not _api_key_for(provider):
        raise ValueError(f"{route_name} route uses {provider!r}, but its API key is not configured.")


def _get_openai_client():
    global _openai_client
    if _openai_client is None:
        from openai import OpenAI
        if config.LLM_PROVIDER.lower() == "groq":
            _openai_client = OpenAI(
                api_key=config.GROQ_API_KEY,
                base_url="https://api.groq.com/openai/v1",
            )
        else:
            _openai_client = OpenAI(api_key=config.OPENAI_API_KEY)
    return _openai_client


def _get_chat_openai_client():
    """Returns the same singleton as _get_openai_client() when providers match."""
    if config.CHAT_LLM_PROVIDER.lower() == config.LLM_PROVIDER.lower():
        return _get_openai_client()
    global _chat_openai_client
    if _chat_openai_client is None:
        from openai import OpenAI
        if config.CHAT_LLM_PROVIDER.lower() == "groq":
            _chat_openai_client = OpenAI(
                api_key=config.GROQ_API_KEY,
                base_url="https://api.groq.com/openai/v1",
            )
        else:
            _chat_openai_client = OpenAI(api_key=config.OPENAI_API_KEY)
    return _chat_openai_client


def _get_anthropic_client():
    global _anthropic_client
    if _anthropic_client is None:
        import anthropic
        _anthropic_client = anthropic.Anthropic(api_key=config.ANTHROPIC_API_KEY)
    return _anthropic_client


def _get_chat_anthropic_client():
    if config.CHAT_LLM_PROVIDER.lower() == config.LLM_PROVIDER.lower():
        return _get_anthropic_client()
    global _chat_anthropic_client
    if _chat_anthropic_client is None:
        import anthropic
        _chat_anthropic_client = anthropic.Anthropic(api_key=config.ANTHROPIC_API_KEY)
    return _chat_anthropic_client


def _get_vision_openai_client():
    global _vision_openai_client
    if _vision_openai_client is None:
        from openai import OpenAI
        if config.VISION_LLM_PROVIDER.lower() == "groq":
            _vision_openai_client = OpenAI(
                api_key=config.GROQ_API_KEY,
                base_url="https://api.groq.com/openai/v1",
            )
        else:
            _vision_openai_client = OpenAI(api_key=config.OPENAI_API_KEY)
    return _vision_openai_client


def _get_vision_anthropic_client():
    global _vision_anthropic_client
    if _vision_anthropic_client is None:
        import anthropic
        _vision_anthropic_client = anthropic.Anthropic(api_key=config.ANTHROPIC_API_KEY)
    return _vision_anthropic_client


def stream_response(
    user_message: str,
    image_base64: str | None = None,
    ambient_context: str = "",
    memory_context: str = "",
    use_tools: bool = False,
) -> Generator[str, None, None]:
    """
    Stream a response from the configured LLM.

    When image_base64 is provided, uses VISION_LLM_PROVIDER/MODEL.
    Otherwise uses LLM_PROVIDER/MODEL.

    Args:
        user_message:     The user's query text.
        image_base64:     Optional base64-encoded PNG for vision input.
        ambient_context:  Plain-text context block (active window, clipboard,
                          focused element) — injected into system prompt.
        memory_context:   Pre-formatted LTM facts + STM session summary from
                          core.memory — injected into system prompt before the
                          ambient context block.
        use_tools:        If True and provider is Anthropic, expose
                  web_search + get_context tools so Claude can
                  pull extra context when it decides to. The model
                  must use the actual tool call interface rather than
                  describing or simulating tool calls in text.
                  Ignored for Groq/OpenAI providers and vision calls.

    Yields:
        Text chunks as they arrive from the API.
    """
    import time
    ts = time.strftime("%H:%M:%S")
    if image_base64:
        print(f"[llm {ts}] User message (vision): {user_message!r}")
    else:
        print(f"[llm {ts}] User message: {user_message!r}")

    if ambient_context:
        _log_context(
            "ambient snapshot — captured at hotkey/voice trigger, injected into system prompt",
            ambient_context,
        )

    if image_base64:
        candidates = _route_candidates(
            config.VISION_LLM_PROVIDER,
            config.VISION_LLM_MODEL,
            config.VISION_LLM_FALLBACKS,
        )
        yield from _stream_with_fallbacks(
            "vision",
            candidates,
            lambda provider, model: _stream_single_response_route(
                provider,
                model,
                user_message,
                image_base64,
                ambient_context,
                memory_context,
                use_tools=False,
                route_name="VISION_LLM",
            ),
        )
    else:
        candidates = _route_candidates(config.LLM_PROVIDER, config.LLM_MODEL, config.LLM_FALLBACKS)
        yield from _stream_with_fallbacks(
            "query",
            candidates,
            lambda provider, model: _stream_single_response_route(
                provider,
                model,
                user_message,
                None,
                ambient_context,
                memory_context,
                use_tools=use_tools,
                route_name="LLM",
            ),
        )


def _stream_with_fallbacks(
    kind: str,
    candidates: list[tuple[str, str]],
    factory,
) -> Generator[str, None, None]:
    last_exc: Exception | None = None
    for idx, (provider, model) in enumerate(candidates):
        emitted = False
        try:
            for chunk in factory(provider, model):
                emitted = True
                yield chunk
            return
        except Exception as exc:
            last_exc = exc
            import time

            ts = time.strftime("%H:%M:%S")
            if emitted:
                print(f"[llm {ts}] Route ({kind}) failed after streaming; not falling back: {exc}")
                raise
            if idx < len(candidates) - 1:
                print(f"[llm {ts}] Route ({kind}) failed before output; trying fallback: {exc}")
                continue
            print(f"[llm {ts}] Route ({kind}) failed; no fallback left: {exc}")
            raise
    if last_exc:
        raise last_exc
    raise ValueError(f"No {kind} model routes configured.")


def _stream_single_response_route(
    provider: str,
    model: str,
    user_message: str,
    image_base64: str | None,
    ambient_context: str,
    memory_context: str,
    use_tools: bool,
    route_name: str,
) -> Generator[str, None, None]:
    _check_route_config(provider, model, route_name)
    effective_model = config.TOOL_LLM_MODEL if provider == "anthropic" and use_tools else model
    _log_model_route("vision" if image_base64 else "query", provider, effective_model, use_tools=use_tools)
    if image_base64:
        if provider in ("groq", "openai"):
            yield from _stream_openai_compat(user_message, image_base64, model, _dynamic_openai_client(provider))
        elif provider == "anthropic":
            yield from _stream_anthropic(user_message, image_base64, model, _dynamic_anthropic_client())
        elif provider == "chatgpt":
            yield from _stream_codex_vision(user_message, image_base64, model, _get_codex_client())
        else:
            raise ValueError(f"Unknown vision provider: {provider}")
        return
    if provider in ("groq", "openai"):
        yield from _stream_openai_compat(user_message, None, model, _dynamic_openai_client(provider), ambient_context, memory_context)
    elif provider == "anthropic":
        yield from _stream_anthropic(user_message, None, effective_model, _dynamic_anthropic_client(), ambient_context, memory_context, use_tools)
    elif provider == "chatgpt":
        yield from _stream_codex(user_message, model, _get_codex_client(), ambient_context, memory_context, use_tools)
    elif provider == "copilot":
        yield from _stream_copilot(user_message, model, ambient_context, memory_context, use_tools)
    else:
        raise ValueError(f"Unknown LLM provider: {provider}")


def _stream_copilot(
    user_message: str,
    model: str,
    ambient_context: str = "",
    memory_context: str = "",
    use_tools: bool = False,
) -> Generator[str, None, None]:
    from core import copilot_client

    parts = []
    if memory_context:
        parts.append(memory_context)
    if ambient_context:
        parts.append("Context:\n" + ambient_context)
    system = "\n\n".join(parts)
    yield from copilot_client.stream(
        user_message,
        model,
        system=system,
        allow_tools=use_tools,
    )


# ------------------------------------------------------------------
# OpenAI / Groq (OpenAI-compatible)
# ------------------------------------------------------------------

def _stream_openai_compat(
    user_message: str,
    image_base64: str | None,
    model: str,
    client,
    ambient_context: str = "",
    memory_context: str = "",
) -> Generator[str, None, None]:
    messages = _build_openai_messages(user_message, image_base64, ambient_context, memory_context)

    with client.chat.completions.create(
        model=model,
        messages=messages,
        stream=True,
        max_tokens=256,
        temperature=0.5,
    ) as stream:
        for chunk in stream:
            delta = chunk.choices[0].delta.content
            if delta:
                yield delta


# ------------------------------------------------------------------
# Codex (ChatGPT subscription) — Responses API streaming
# ------------------------------------------------------------------

def _build_codex_text(user_message: str, ambient_context: str = "", memory_context: str = "") -> str:
    """Prepend context into a single text block for the Codex endpoint."""
    parts = []
    if memory_context:
        parts.append(memory_context)
    if ambient_context:
        parts.append(f"---\n{ambient_context}")
    parts.append(user_message)
    return "\n\n".join(parts)


def _stream_codex(
    user_message: str,
    model: str,
    client,
    ambient_context: str = "",
    memory_context: str = "",
    use_tools: bool = False,
) -> Generator[str, None, None]:
    """Stream a response via the Codex endpoint using the Responses API."""
    # The Responses API tool-call loop (previous_response_id chaining) is not
    # implemented here.  When tools are requested we eagerly inject context that
    # Claude would otherwise fetch on demand: open supported documents and clipboard.
    if use_tools:
        doc_text = read_active_document_for_context()
        if doc_text:
            doc_block = f"[Open documents]\n{doc_text}"
            ambient_context = f"{ambient_context}\n\n---\n{doc_block}".strip() if ambient_context else doc_block
    text = _build_codex_text(user_message, ambient_context, memory_context)
    with client.responses.stream(
        model=model,
        input=[{"type": "message", "role": "user", "content": [{"type": "input_text", "text": text}]}],
        instructions=config.get_system_prompt(),
        store=False,
    ) as stream:
        for event in stream:
            if getattr(event, 'type', '') == 'response.output_text.delta':
                delta = getattr(event, 'delta', '')
                if delta:
                    yield delta


def _stream_codex_vision(
    user_message: str,
    image_base64: str,
    model: str,
    client,
) -> Generator[str, None, None]:
    """Stream a vision response via the Codex endpoint (Responses API with image input)."""
    input_content = [
        {"type": "input_text",  "text": user_message},
        {"type": "input_image", "image_url": f"data:image/png;base64,{image_base64}"},
    ]
    with client.responses.stream(
        model=model,
        input=[{"type": "message", "role": "user", "content": input_content}],
        instructions=config.get_system_prompt(),
        store=False,
    ) as stream:
        for event in stream:
            if getattr(event, 'type', '') == 'response.output_text.delta':
                delta = getattr(event, 'delta', '')
                if delta:
                    yield delta


def _build_openai_messages(
    user_message: str,
    image_base64: str | None,
    ambient_context: str = "",
    memory_context: str = "",
) -> list:
    system = config.get_system_prompt()
    if memory_context:
        system += f"\n\n{memory_context}"
    if ambient_context:
        system += f"\n\n---\n{ambient_context}"
    if image_base64:
        content = [
            {"type": "text", "text": user_message},
            {
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{image_base64}"},
            },
        ]
    else:
        content = user_message

    return [
        {"role": "system", "content": system},
        {"role": "user", "content": content},
    ]


# ------------------------------------------------------------------
# Anthropic Claude  —  shared tool-loop helper
# ------------------------------------------------------------------

def _run_anthropic_tool_loop(
    client,
    messages: list,
    first_response,
    model: str,
    system: str,
    max_tokens: int,
) -> Generator[str, None, None]:
    """
    Execute Anthropic tool calls and yield text from subsequent rounds.
    *messages* is mutated in place (assistant + tool-result turns are appended).
    """
    messages.append({"role": "assistant", "content": first_response.content})
    final = first_response
    for _round in range(3):
        tool_results = []
        for block in final.content:
            if getattr(block, "type", "") == "tool_use":
                result = _execute_model_tool(block.name, block.input or {})
                tool_results.append({
                    "type": "tool_result",
                    "tool_use_id": block.id,
                    "content": result,
                })
        if not tool_results:
            return
        messages.append({"role": "user", "content": tool_results})
        response = client.messages.create(
            model=model,
            max_tokens=max_tokens,
            system=system,
            messages=messages,
            tools=_get_tool_schemas(),
        )
        for block in response.content:
            if getattr(block, "type", "") == "text":
                text = getattr(block, "text", "")
                if text:
                    yield text
        if response.stop_reason != "tool_use":
            return
        final = response
        messages.append({"role": "assistant", "content": response.content})


def _stream_anthropic(
    user_message: str,
    image_base64: str | None,
    model: str,
    client,
    ambient_context: str = "",
    memory_context: str = "",
    use_tools: bool = False,
) -> Generator[str, None, None]:
    system = config.get_system_prompt()
    if memory_context:
        system += f"\n\n{memory_context}"
    if ambient_context:
        system += f"\n\n---\n{ambient_context}"

    if image_base64:
        content = [
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/png",
                    "data": image_base64,
                },
            },
            {"type": "text", "text": user_message},
        ]
    else:
        content = user_message

    # --- No tools: original streaming path (lowest latency) ---
    if not use_tools:
        with client.messages.stream(
            model=model,
            max_tokens=256,
            system=system,
            messages=[{"role": "user", "content": content}],
        ) as stream:
            for text in stream.text_stream:
                yield text
        return

    # --- Tool-enabled path: stream first round for fast first-token ---
    # If no tool is called (common case), text streams immediately.
    # Only falls back to blocking create() if Claude actually invokes a tool.
    messages: list[dict] = [{"role": "user", "content": content}]

    with client.messages.stream(
        model=model,
        max_tokens=512,
        system=system,
        messages=messages,
        tools=_get_tool_schemas(),
    ) as stream:
        for text in stream.text_stream:
            yield text
        final = stream.get_final_message()

    if final.stop_reason != "tool_use":
        return

    # A tool was called — execute it and do followup round(s) non-streaming.
    yield from _run_anthropic_tool_loop(client, messages, final, model, system, max_tokens=512)


# ------------------------------------------------------------------
# Inline rewrite / fix  (Ctrl+Shift+Q)
# ------------------------------------------------------------------

_REWRITE_SYSTEM_PROMPT = (
    "You are a text editor assistant. "
    "Rewrite or fix the provided text. "
    "Output ONLY the corrected/rewritten text — no explanation, "
    "no markdown, no commentary, no code fences. "
    "Your entire response will be pasted directly as a replacement for the original text."
)


def stream_rewrite(selected_text: str, intent_prompt: str = "Rewrite or fix the following text") -> Generator[str, None, None]:
    """
    Stream a rewrite/fix of the selected text using the primary LLM.

    The system prompt instructs the model to output raw replacement text only —
    no prose, no markdown, no explanation.  The result is pasted back directly.

    Args:
        selected_text:  The text to rewrite.
        intent_prompt:  The instruction (e.g. "Fix the grammar and spelling").
                        Taken from the caller's chosen intent row.
    """
    import time
    ts = time.strftime("%H:%M:%S")
    print(f"[llm {ts}] Rewrite request ({len(selected_text)} chars) — {intent_prompt[:60]!r}")
    user_message = f"{intent_prompt}:\n\n{selected_text}"
    candidates = _route_candidates(config.LLM_PROVIDER, config.LLM_MODEL, config.LLM_FALLBACKS)
    yield from _stream_with_fallbacks(
        "rewrite",
        candidates,
        lambda provider, model: _stream_single_rewrite_route(provider, model, user_message),
    )


def _stream_single_rewrite_route(provider: str, model: str, user_message: str) -> Generator[str, None, None]:
    _check_route_config(provider, model, "LLM")
    _log_model_route("rewrite", provider, model, use_tools=False)
    if provider in ("groq", "openai"):
        client = _dynamic_openai_client(provider)
        with client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": _REWRITE_SYSTEM_PROMPT},
                {"role": "user",   "content": user_message},
            ],
            stream=True,
            max_tokens=1024,
            temperature=0.3,
        ) as stream:
            for chunk in stream:
                delta = chunk.choices[0].delta.content
                if delta:
                    yield delta
    elif provider == "anthropic":
        client = _dynamic_anthropic_client()
        with client.messages.stream(
            model=model,
            max_tokens=1024,
            system=_REWRITE_SYSTEM_PROMPT,
            messages=[{"role": "user", "content": user_message}],
            temperature=0.3,
        ) as stream:
            for text in stream.text_stream:
                yield text
    elif provider == "copilot":
        from core import copilot_client
        yield from copilot_client.stream(
            user_message,
            model,
            system=_REWRITE_SYSTEM_PROMPT,
        )
    elif provider == "chatgpt":
        with _get_codex_client().responses.stream(
            model=model,
            input=[{"type": "message", "role": "user", "content": [{"type": "input_text", "text": user_message}]}],
            instructions=_REWRITE_SYSTEM_PROMPT,
            store=False,
        ) as stream:
            for event in stream:
                if getattr(event, 'type', '') == 'response.output_text.delta':
                    delta = getattr(event, 'delta', '')
                    if delta:
                        yield delta
    else:
        raise ValueError(f"Unknown rewrite provider: {provider}")


# ------------------------------------------------------------------
# Multi-turn (chat window)
# ------------------------------------------------------------------

def stream_response_with_history(
    messages: list,
    memory_context: str = "",
) -> Generator[str, None, None]:
    """
    Stream a response given a pre-built messages list including history.
    Uses CHAT_LLM_PROVIDER / CHAT_LLM_MODEL (defaults to LLM_PROVIDER / LLM_MODEL).

    Args:
        messages:        [{{"role": "system"|"user"|"assistant", "content": str}}, ...]
        memory_context:  Pre-formatted LTM facts from core.memory — appended to
                         the system message so the model is aware of user facts.
    """
    # Inject memory context into the system message (or prepend one)
    if memory_context:
        sys_idx = next(
            (i for i, m in enumerate(messages) if m["role"] == "system"), None
        )
        if sys_idx is not None:
            messages = list(messages)   # shallow copy — don't mutate the caller's list
            messages[sys_idx] = {
                **messages[sys_idx],
                "content": messages[sys_idx]["content"] + f"\n\n{memory_context}",
            }
        else:
            messages = [{"role": "system", "content": memory_context}] + list(messages)
    candidates = _route_candidates(
        config.CHAT_LLM_PROVIDER,
        config.CHAT_LLM_MODEL,
        config.CHAT_LLM_FALLBACKS,
    )
    yield from _stream_with_fallbacks(
        "chat",
        candidates,
        lambda provider, model: _stream_single_history_route(provider, model, messages),
    )


def _stream_single_history_route(provider: str, model: str, messages: list) -> Generator[str, None, None]:
    _check_route_config(provider, model, "CHAT_LLM")
    _log_model_route("chat", provider, model, use_tools=(provider == "anthropic"))
    if provider in ("groq", "openai"):
        client = _dynamic_openai_client(provider)
        with client.chat.completions.create(
            model=model,
            messages=messages,
            stream=True,
            max_tokens=1024,
            temperature=0.7,
        ) as stream:
            for chunk in stream:
                delta = chunk.choices[0].delta.content
                if delta:
                    yield delta
    elif provider == "anthropic":
        client = _dynamic_anthropic_client()
        system = next((m["content"] for m in messages if m["role"] == "system"), "")
        _VALID_KEYS = {"role", "content"}
        turns = [
            {k: v for k, v in m.items() if k in _VALID_KEYS}
            for m in messages if m["role"] != "system"
        ]
        # Tool-enabled loop — stream first round, block only if a tool is actually called.
        with client.messages.stream(
            model=model,
            max_tokens=1024,
            system=system,
            messages=turns,
            tools=_get_tool_schemas(),
        ) as stream:
            for text in stream.text_stream:
                yield text
            final = stream.get_final_message()

        if final.stop_reason != "tool_use":
            return

        yield from _run_anthropic_tool_loop(client, turns, final, model, system, max_tokens=1024)
    elif provider == "chatgpt":
        # For the chat history path, extract the last user message and use Responses API.
        # The full history is packed into a single user turn with prior turns prefixed.
        system_msg = next((m["content"] for m in messages if m["role"] == "system"), "")
        turns = [m for m in messages if m["role"] != "system"]
        # Build a text representation of prior context for the model
        last_user = turns[-1]["content"] if turns else ""
        history_prefix = ""
        for m in turns[:-1]:
            label = "User" if m["role"] == "user" else "Assistant"
            history_prefix += f"{label}: {m['content']}\n"
        full_input = (history_prefix + last_user) if history_prefix else last_user
        with _get_chat_codex_client().responses.stream(
            model=model,
            input=[{"type": "message", "role": "user", "content": [{"type": "input_text", "text": full_input}]}],
            instructions=system_msg,
            store=False,
        ) as stream:
            for event in stream:
                if getattr(event, 'type', '') == 'response.output_text.delta':
                    delta = getattr(event, 'delta', '')
                    if delta:
                        yield delta
    elif provider == "copilot":
        from core import copilot_client

        system_msg = next((m["content"] for m in messages if m["role"] == "system"), "")
        turns = [m for m in messages if m["role"] != "system"]
        full_input = ""
        for m in turns:
            label = "User" if m["role"] == "user" else "Assistant"
            full_input += f"{label}: {m['content']}\n"
        yield from copilot_client.stream(
            full_input.strip(),
            model,
            system=system_msg,
            session_id="wisp-chat",
        )
    else:
        raise ValueError(f"Unknown chat LLM provider: {provider}")
